import pandas as pd
from pptx import Presentation
import copy
import math

excel_file = "BADGE (CLEAN FILE).xlsx"
template_file = "CAMPUSMINISTRYBADGE.pptx"
output_file = "filled_badges.pptx"

def assign_tables_and_discussions(df):
    """
    Assign table and discussion numbers based on language preferences
    35 tables for dorm seating (8 people per table)
    35 tables for discussion (10 people per table)
    """
    
    # Separate by language preference
    amharic_speakers = df[df['What is your preferred language for sermons and/or engaging in group discussions?'] == 'Amharic'].copy()
    english_speakers = df[df['What is your preferred language for sermons and/or engaging in group discussions?'] == 'English'].copy()
    no_preference = df[df['What is your preferred language for sermons and/or engaging in group discussions?'].isna()].copy()
    
    # Treat no preference as English speakers
    english_speakers = pd.concat([english_speakers, no_preference])
    
    print(f"Amharic speakers: {len(amharic_speakers)}")
    print(f"English speakers: {len(english_speakers)}")
    
    # Assign discussion tables (10 people per table)
    discussion_table = 1
    discussion_count = 0
    
    # Assign Amharic speakers to discussion tables first
    for idx in amharic_speakers.index:
        if discussion_count == 10:
            discussion_table += 1
            discussion_count = 0
        df.loc[idx, 'DISCUSSION #'] = discussion_table
        discussion_count += 1
    
    # Continue with English speakers
    if discussion_count == 10:
        discussion_table += 1
        discussion_count = 0
        
    for idx in english_speakers.index:
        if discussion_count == 10:
            discussion_table += 1
            discussion_count = 0
        df.loc[idx, 'DISCUSSION #'] = discussion_table
        discussion_count += 1
    
    # Assign table numbers (8 people per table)
    table_number = 1
    table_count = 0
    
    # Assign Amharic speakers to tables (with DA prefix for Amharic)
    for idx in amharic_speakers.index:
        if table_count == 8:
            table_number += 1
            table_count = 0
        df.loc[idx, 'TABLE #'] = f"DA {table_number:02d}"
        table_count += 1
    
    # Continue with English speakers (with regular numbers)
    if table_count == 8:
        table_number += 1
        table_count = 0
        
    for idx in english_speakers.index:
        if table_count == 8:
            table_number += 1
            table_count = 0
        df.loc[idx, 'TABLE #'] = f"{table_number:02d}"
        table_count += 1
    
    return df

# Load and process data
df = pd.read_excel(excel_file)
df = assign_tables_and_discussions(df)

valid_rows = []
skipped_rows = []

for index, row in df.iterrows():
    missing_fields = []
    
    if pd.isna(row['First Name']) or str(row['First Name']).strip() == '':
        missing_fields.append('First Name')
    if pd.isna(row['Last Name']) or str(row['Last Name']).strip() == '':
        missing_fields.append('Last Name')
    
    if missing_fields:
        skipped_rows.append(f"Row {index + 2}: Missing {', '.join(missing_fields)}")
    else:
        full_name = f"{str(row['First Name']).strip().upper()} {str(row['Last Name']).strip().upper()}"
        
        campus = str(row['Which campus ministry are you a part of? In case, you are not part of the listed campus ministries, it\'s open to you as well, and please come and join us!']) if not pd.isna(row['Which campus ministry are you a part of? In case, you are not part of the listed campus ministries, it\'s open to you as well, and please come and join us!']) else 'NEW'
        if 'I don\'t have' in campus or campus == 'nan':
            campus = 'NEW'
        
        role = 'PARTICIPANT'
        dorm = str(row['DORM NUMBER']).strip() if not pd.isna(row['DORM NUMBER']) else ''
        table = str(row['TABLE #']).strip() if not pd.isna(row['TABLE #']) else ''
        discussion = str(row['DISCUSSION #']).strip() if not pd.isna(row['DISCUSSION #']) else ''
        
        valid_rows.append({
            'Full_Name': full_name,
            'Campus': campus,
            'Role': role,
            'Dorm': dorm,
            'Table': table,
            'Discussion': discussion
        })

print(f"Processing {len(valid_rows)} valid rows")
if skipped_rows:
    print("Skipped rows:")
    for skip in skipped_rows:
        print(f"  {skip}")

template_pres = Presentation(template_file)
template_slide = template_pres.slides[0]

def duplicate_slide(pres, source_slide):
    try:
        blank_slide_layout = pres.slide_layouts[6]
    except:
        blank_slide_layout = pres.slide_layouts[0]
    
    copied_slide = pres.slides.add_slide(blank_slide_layout)
    
    for shape in source_slide.shapes:
        try:
            el = shape.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        except Exception as e:
            print(f"Warning: Could not copy shape: {e}")
            continue
    
    try:
        for rel_id, rel in source_slide.part.rels.items():
            if "notesSlide" not in rel.reltype:
                copied_slide.part.rels.add_relationship(
                    rel.reltype, rel._target, rel_id
                )
    except Exception as e:
        print(f"Warning: Could not copy relationships: {e}")
    
    return copied_slide

output_pres = Presentation()
output_pres.slide_width = template_pres.slide_width
output_pres.slide_height = template_pres.slide_height

num_batches = math.ceil(len(valid_rows) / 4)

for batch_idx in range(num_batches):
    start_idx = batch_idx * 4
    end_idx = min(start_idx + 4, len(valid_rows))
    batch_data = valid_rows[start_idx:end_idx]
    
    new_slide = duplicate_slide(output_pres, template_slide)

    for i, row_data in enumerate(batch_data):
        badge_num = i + 1
        
        name_placeholder = f"{{{{NAME{badge_num}}}}}"
        campus_placeholder = f"{{{{CAMPUS{badge_num}}}}}"
        role_placeholder = f"{{{{ROLE{badge_num}}}}}"
        dorm_placeholder = f"{{{{DORM{badge_num}}}}}"
        table_placeholder = f"{{{{TABLE{badge_num}}}}}"
        discussion_placeholder = f"{{{{DISCUSSION{badge_num}}}}}"
        
        for shape in new_slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if name_placeholder in run.text:
                            run.text = run.text.replace(name_placeholder, row_data['Full_Name'])
                            if len(row_data['Full_Name']) > 20:
                                if run.font.size:
                                    run.font.size = int(run.font.size * 0.7)
                            elif len(row_data['Full_Name']) > 15:
                                if run.font.size:
                                    run.font.size = int(run.font.size * 0.8)
                        if campus_placeholder in run.text:
                            run.text = run.text.replace(campus_placeholder, row_data['Campus'])
                            if len(row_data['Campus']) > 25:
                                if run.font.size:
                                    run.font.size = int(run.font.size * 0.7)
                            elif len(row_data['Campus']) > 20:
                                if run.font.size:
                                    run.font.size = int(run.font.size * 0.8)
                        if role_placeholder in run.text:
                            run.text = run.text.replace(role_placeholder, row_data['Role'])
                        if dorm_placeholder in run.text:
                            run.text = run.text.replace(dorm_placeholder, row_data['Dorm'])
                        if table_placeholder in run.text:
                            run.text = run.text.replace(table_placeholder, row_data['Table'])
                        if discussion_placeholder in run.text:
                            run.text = run.text.replace(discussion_placeholder, row_data['Discussion'])

    remaining_badges = 4 - len(batch_data)
    for i in range(remaining_badges):
        badge_num = len(batch_data) + i + 1
        
        name_placeholder = f"{{{{NAME{badge_num}}}}}"
        campus_placeholder = f"{{{{CAMPUS{badge_num}}}}}"
        role_placeholder = f"{{{{ROLE{badge_num}}}}}"
        dorm_placeholder = f"{{{{DORM{badge_num}}}}}"
        table_placeholder = f"{{{{TABLE{badge_num}}}}}"
        discussion_placeholder = f"{{{{DISCUSSION{badge_num}}}}}"
        
        for shape in new_slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace(name_placeholder, "")
                        run.text = run.text.replace(campus_placeholder, "")
                        run.text = run.text.replace(role_placeholder, "")
                        run.text = run.text.replace(dorm_placeholder, "")
                        run.text = run.text.replace(table_placeholder, "")
                        run.text = run.text.replace(discussion_placeholder, "")

output_pres.save(output_file)
print(f"Saved {num_batches} slides to {output_file}")

# Save the updated Excel file with table assignments
df.to_excel("BADGE_WITH_ASSIGNMENTS.xlsx", index=False)
print("Saved updated Excel file with table assignments: BADGE_WITH_ASSIGNMENTS.xlsx")