import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
import copy
import math
import os

# Get the absolute path to the project root directory
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))

excel_file = os.path.join(project_root, "data", "excel", "BADGE (CLEAN FILE).xlsx")
excel_file_2 = os.path.join(project_root, "data", "excel", "General Assembly Dorm Assignment1.xlsx")
excel_file_3 = os.path.join(project_root, "data", "excel", "General Assembly Dorm Assignment2.xlsx")
template_file = os.path.join(project_root, "templates", "CAMPUSMINISTRYBADGE_PATCHED.pptx")
output_file = os.path.join(project_root, "output", "badges", "filled_badges.pptx")

campus_ministry_locations = {
    # Atlanta variations
    "Atlanta": "ATLANTA CAMPUS MINISTRY",
    "Atlanta, Georgia Campus Ministry": "ATLANTA CAMPUS MINISTRY",
    "Atlanta Coordinator": "ATLANTA CAMPUS MINISTRY",
    
    # DC/Washington variations
    "DC": "GEORGE WASHINGTON UNIVERSITY CAMPUS MINISTRY",
    "Dc": "GEORGE WASHINGTON UNIVERSITY CAMPUS MINISTRY",
    "DC campus Ministry": "GEORGE WASHINGTON UNIVERSITY CAMPUS MINISTRY",
    "George Washington University Campus Ministry": "GEORGE WASHINGTON UNIVERSITY CAMPUS MINISTRY",
    
    # Virginia variations
    "Virginia": "GEORGE MASON UNIVERSITY CAMPUS MINISTRY",
    "Verginia": "GEORGE MASON UNIVERSITY CAMPUS MINISTRY",
    "Verginia ": "GEORGE MASON UNIVERSITY CAMPUS MINISTRY",
    "Virginia ": "GEORGE MASON UNIVERSITY CAMPUS MINISTRY",
    "VA Campus Ministry": "GEORGE MASON UNIVERSITY CAMPUS MINISTRY",
    "George Mason University (GMU) Campus Ministry": "GEORGE MASON UNIVERSITY CAMPUS MINISTRY",
    "Ve": "GEORGE MASON UNIVERSITY CAMPUS MINISTRY",
    "Ver": "GEORGE MASON UNIVERSITY CAMPUS MINISTRY",
    
    # NOVA variations
    "Verginia North": "NOVA CAMPUS MINISTRY",
    "Nova Campus Ministry": "NOVA CAMPUS MINISTRY",
    
    # Boston variations
    "Boston": "BOSTON CAMPUS MINISTRY",
    "Bosten": "BOSTON CAMPUS MINISTRY",
    "Bosten ": "BOSTON CAMPUS MINISTRY",
    "Boston Campus Ministry": "BOSTON CAMPUS MINISTRY",
    
    # New York variations
    "New York": "NEW YORK CAMPUS MINISTRY",
    "Newyork": "NEW YORK CAMPUS MINISTRY",
    "Newyork ": "NEW YORK CAMPUS MINISTRY",
    
    # Carolina variations
    "Carolina": "NORTH CAROLINA CAMPUS MINISTRY",
    " N Carolina ": "NORTH CAROLINA CAMPUS MINISTRY",
    
    # Indiana variations
    "Indiana": "INDIANA CAMPUS MINISTRY",
    "Indiana ": "INDIANA CAMPUS MINISTRY",
    
    # Ohio variations
    "Ohio": "OHIO CAMPUS MINISTRY",
    "Ohio ": "OHIO CAMPUS MINISTRY",
    
    # Los Angeles variations
    "Los Angeles": "LOS ANGELES CAMPUS MINISTRY",
    "Los Angelos Campus Ministry": "LOS ANGELES CAMPUS MINISTRY",
    
    # Denver variations
    "Denver": "DENVER CAMPUS MINISTRY",
    "Denver Campus Ministry": "DENVER CAMPUS MINISTRY",
    "Denver Coordinator": "DENVER CAMPUS MINISTRY",
    
    # Seattle variations
    "Seattle": "SEATTLE CAMPUS MINISTRY",
    "Seattle Campus Ministry": "SEATTLE CAMPUS MINISTRY",
    "Seattle Coordinator": "SEATTLE CAMPUS MINISTRY",
    
    # Nashville variations
    "Nashville": "NASHVILLE CAMPUS MINISTRY",
    "Nashville Campus Ministry": "NASHVILLE CAMPUS MINISTRY",
    
    # Minnesota variations
    "Minnesota": "MINNESOTA CAMPUS MINISTRY",
    "Minnesota Campus Ministry": "MINNESOTA CAMPUS MINISTRY",
    
    # San Jose variations
    "San Jose": "SAN JOSE CAMPUS MINISTRY",
    "San Jose Campus Ministry": "SAN JOSE CAMPUS MINISTRY",
    
    # UMD variations
    "UMD": "UNIVERSITY OF MARYLAND (UMD) CAMPUS MINISTRY",
    "University of Maryland (UMD) Campus Ministry": "UNIVERSITY OF MARYLAND (UMD) CAMPUS MINISTRY",
    
    # Virginia Tech variations
    "VTech": "VIRGINIA TECH CAMPUS MINISTRY",
    "Virginia Tech Campus Ministry(VTech)": "VIRGINIA TECH CAMPUS MINISTRY",
    
    # Virtual Campus Ministry variations
    "Virtual": "VIRTUAL CAMPUS MINISTRY (VCM)",
    "Virtual campus Ministry(VCM)": "VIRTUAL CAMPUS MINISTRY (VCM)",
    "Virtual campus Ministry": "VIRTUAL CAMPUS MINISTRY (VCM)",
    "Virtual Campus Ministry": "VIRTUAL CAMPUS MINISTRY (VCM)",
    "VCM": "VIRTUAL CAMPUS MINISTRY (VCM)",
    "VCM Coordinator": "VIRTUAL CAMPUS MINISTRY (VCM)",
    
    # MIU variations
    "MIU": "MAHARISHI INTERNATIONAL UNIVERSITY (MIU) CAMPUS MINISTRY",
    "Maharishi International University (MIU) Campus Ministry": "MAHARISHI INTERNATIONAL UNIVERSITY (MIU) CAMPUS MINISTRY",
    
    # Montgomery variations
    "Montgomery": "MONTGOMERY COLLEGE CAMPUS MINISTRY",
    "Montgomery College Campus Ministry(MC Campus Ministry)": "MONTGOMERY COLLEGE CAMPUS MINISTRY",
    
    # Special cases
    "Guest": "GUEST",
    "NEW": "NEW",
    "New": "NEW",
    "Graduated": "GRADUATED FROM CAMPUS MINISTRY",
    "Graduated from Campus Ministry": "GRADUATED FROM CAMPUS MINISTRY",
    "I don't have Campus Ministry": "NO CAMPUS MINISTRY"
}

def extract_language_preference(first_name):
    """
    Extract language preference from first name field
    Examples: 
    - 'Solyana E' -> 'English'
    - 'Ruth -A' -> 'Amharic' 
    - 'NAHOM&A' -> 'Amharic'
    - 'SARA&E' -> 'English'
    - 'Miriam E&A' -> 'English' (E takes precedence)
    """
    if pd.isna(first_name) or str(first_name).strip() == '':
        return 'English'
    
    name_str = str(first_name).strip().upper()
    
    # Check for Amharic indicators (in order of priority)
    if '&A' in name_str and '&E' not in name_str:
        return 'Amharic'
    elif ' -A' in name_str or '- A' in name_str or '-A' in name_str:
        return 'Amharic'
    
    # Check for English indicators
    if '&E' in name_str:
        return 'English'
    elif ' E' in name_str or name_str.endswith('E'):
        return 'English'
    
    # Everything else defaults to English (including New, etc.)
    return 'English'

def process_second_excel_file(file_path):
    """
    Process the second Excel file and return a DataFrame with the same structure as the first
    Only include people with "1" at the beginning of their first name
    """
    df2 = pd.read_excel(file_path)
    
    # Filter out rows with NaN first names
    df2 = df2.dropna(subset=['First Name'])
    
    # Create a new DataFrame with the same structure as the first Excel file
    processed_data = []
    
    for index, row in df2.iterrows():
        # Extract clean first name and language preference
        first_name_raw = str(row['First Name']).strip()
        
        # Only process if first name starts with "1"
        if not first_name_raw.startswith('1'):
            continue
        
        language_pref = extract_language_preference(first_name_raw)
        
        # Clean the first name (remove "1" and language indicators)
        clean_first_name = first_name_raw[1:]  # Remove the "1" prefix
        # Use comprehensive cleaning function
        clean_first_name = clean_name_comprehensive(clean_first_name)
        
        # Get campus ministry
        campus_ministry = str(row['Campus Minstry']) if not pd.isna(row['Campus Minstry']) else 'NEW'
        if 'I don\'t have' in campus_ministry or campus_ministry == 'nan':
            campus_ministry = 'NEW'
        
        # Create dorm info from Floor and Room - fix the extra digit issue
        floor = str(row['Floor ']).strip() if not pd.isna(row['Floor ']) else ''
        room = str(row['Room']).strip() if not pd.isna(row['Room']) else ''
        
        # Fix dorm number - remove extra digit (3312C -> 312C)
        if floor and room:
            dorm_raw = f"{floor}{room}"
            # If dorm starts with digit and is 4+ characters, remove first digit
            if len(dorm_raw) >= 4 and dorm_raw[0].isdigit() and dorm_raw[1:4].isdigit():
                dorm = dorm_raw[1:]
            else:
                dorm = dorm_raw
        else:
            dorm = 'N/A'
        
        processed_data.append({
            'First Name ': clean_first_name,
            'Last Name': str(row['Last Name']).strip() if not pd.isna(row['Last Name']) else '',
            'Which campus ministry are you a part of? In case, you are not part of the listed campus ministries, it\'s open to you as well, and please come and join us! ': campus_ministry,
            'What is your preferred language for sermons and/or engaging in group discussions?': language_pref,
            'DORM NUMBER': dorm,
            'TABLE #': '',  # Will be assigned later
            'DISCUSSION #': ''  # Will be assigned later
        })
    
    return pd.DataFrame(processed_data)

def assign_tables_and_discussions(df):
    """
    Assign table and discussion numbers based on language preferences
    Max 35 tables for discussion (10 people each)
    Max 35 tables for seating (8 people each)
    Systematically distribute 1 person from each campus ministry per table for maximum diversity
    """
    amharic_speakers = df[df['What is your preferred language for sermons and/or engaging in group discussions?'] == 'Amharic'].copy()
    english_speakers = df[df['What is your preferred language for sermons and/or engaging in group discussions?'] == 'English'].copy()
    no_preference = df[df['What is your preferred language for sermons and/or engaging in group discussions?'].isna()].copy()
    
    english_speakers = pd.concat([english_speakers, no_preference])
    
    print(f"Amharic speakers: {len(amharic_speakers)}")
    print(f"English speakers: {len(english_speakers)}")
    
    def distribute_by_ministry(people_df, table_prefix="", table_size=8):
        """Distribute people ensuring 1 person from each ministry per table when possible"""
        # Group by campus ministry
        ministry_groups = {}
        for idx, row in people_df.iterrows():
            ministry = row['Which campus ministry are you a part of? In case, you are not part of the listed campus ministries, it\'s open to you as well, and please come and join us! ']
            if pd.isna(ministry):
                ministry = 'NEW'
            if ministry not in ministry_groups:
                ministry_groups[ministry] = []
            ministry_groups[ministry].append(idx)
        
        # Shuffle each ministry group
        for ministry in ministry_groups:
            import random
            random.shuffle(ministry_groups[ministry])
        
        print(f"  Ministry groups: {[(k, len(v)) for k, v in ministry_groups.items()]}")
        
        # Distribute people
        table_assignments = {}
        table_number = 1
        table_counts = {}
        
        # Phase 1: Distribute one person from each ministry to different tables
        ministry_names = list(ministry_groups.keys())
        ministry_pointers = {m: 0 for m in ministry_names}
        
        while any(ministry_pointers[m] < len(ministry_groups[m]) for m in ministry_names):
            for ministry in ministry_names:
                if ministry_pointers[ministry] < len(ministry_groups[ministry]):
                    if table_number not in table_counts:
                        table_counts[table_number] = 0
                    
                    if table_counts[table_number] < table_size and table_number <= 35:
                        person_idx = ministry_groups[ministry][ministry_pointers[ministry]]
                        table_assignments[person_idx] = table_number
                        table_counts[table_number] += 1
                        ministry_pointers[ministry] += 1
                        
                        if table_counts[table_number] >= table_size:
                            table_number += 1
                    else:
                        # Move to next table if current is full
                        table_number += 1
                        if table_number > 35:
                            # Put overflow in table 35
                            person_idx = ministry_groups[ministry][ministry_pointers[ministry]]
                            table_assignments[person_idx] = 35
                            ministry_pointers[ministry] += 1
        
        return table_assignments, table_number
    
    def distribute_discussions(people_df, discussion_size=10):
        """Distribute discussion groups with diversity"""
        # Group by campus ministry
        ministry_groups = {}
        for idx, row in people_df.iterrows():
            ministry = row['Which campus ministry are you a part of? In case, you are not part of the listed campus ministries, it\'s open to you as well, and please come and join us! ']
            if pd.isna(ministry):
                ministry = 'NEW'
            if ministry not in ministry_groups:
                ministry_groups[ministry] = []
            ministry_groups[ministry].append(idx)
        
        # Shuffle each ministry group
        for ministry in ministry_groups:
            import random
            random.shuffle(ministry_groups[ministry])
        
        # Distribute people
        discussion_assignments = {}
        discussion_number = 1
        discussion_counts = {}
        
        # Phase 1: Distribute one person from each ministry to different discussions
        ministry_names = list(ministry_groups.keys())
        ministry_pointers = {m: 0 for m in ministry_names}
        
        while any(ministry_pointers[m] < len(ministry_groups[m]) for m in ministry_names):
            for ministry in ministry_names:
                if ministry_pointers[ministry] < len(ministry_groups[ministry]):
                    if discussion_number not in discussion_counts:
                        discussion_counts[discussion_number] = 0
                    
                    if discussion_counts[discussion_number] < discussion_size and discussion_number <= 35:
                        person_idx = ministry_groups[ministry][ministry_pointers[ministry]]
                        discussion_assignments[person_idx] = discussion_number
                        discussion_counts[discussion_number] += 1
                        ministry_pointers[ministry] += 1
                        
                        if discussion_counts[discussion_number] >= discussion_size:
                            discussion_number += 1
                    else:
                        # Move to next discussion if current is full
                        discussion_number += 1
                        if discussion_number > 35:
                            break
        
        return discussion_assignments, discussion_number
    
    # Assign Amharic speakers to tables
    amharic_table_assignments, amharic_max_table = distribute_by_ministry(amharic_speakers, "DA", 8)
    
    # Assign English speakers to tables (continuing from where Amharic left off)
    english_table_assignments, english_max_table = distribute_by_ministry(english_speakers, "", 8)
    
    # Apply table assignments with T prefix (Saturday only, Sunday will be randomized during badge generation)
    for idx, table_num in amharic_table_assignments.items():
        sat_table = f"T{table_num}"
        df.loc[idx, 'TABLE #'] = sat_table
    
    for idx, table_num in english_table_assignments.items():
        sat_table = f"T{table_num}"
        df.loc[idx, 'TABLE #'] = sat_table
    
    # Assign discussions
    amharic_discussion_assignments, amharic_max_discussion = distribute_discussions(amharic_speakers, 10)
    english_discussion_assignments, english_max_discussion = distribute_discussions(english_speakers, 10)
    
    # Apply discussion assignments with language prefix
    for idx, discussion_num in amharic_discussion_assignments.items():
        df.loc[idx, 'DISCUSSION #'] = f"DA {discussion_num}"
    
    for idx, discussion_num in english_discussion_assignments.items():
        df.loc[idx, 'DISCUSSION #'] = f"DE {discussion_num}"
    
    print(f"Final discussion table count: {max(amharic_max_discussion, english_max_discussion)}")
    print(f"Final seating table count: {max(amharic_max_table, english_max_table)}")
    
    return df

def clean_name_comprehensive(name_str):
    """
    Comprehensive name cleaning function that removes all language indicators
    """
    if pd.isna(name_str) or str(name_str).strip() == '':
        return ''
    
    clean_name = str(name_str).strip()
    
    # Remove all possible language indicators comprehensively
    for indicator in ['&A', '&E', ' &A', ' &E', '&A ', '&E ', ' -A', '- A', '-A', ' -E', '- E', '-E', ' E&A', 'E&A', ' E', 'E', ' A', 'A', ' -New', '- New', '-New', ' New', 'New']:
        clean_name = clean_name.replace(indicator, '')
    
    return clean_name.strip()

def process_main_excel_file(df):
    """
    Process the main Excel file to handle names with &A and &E language indicators
    Clean names and update language preferences if needed
    """
    for index, row in df.iterrows():
        first_name_raw = str(row['First Name ']).strip() if not pd.isna(row['First Name ']) else ''
        
        if first_name_raw and ('&A' in first_name_raw or '&E' in first_name_raw):
            # Extract language preference from name
            language_pref = extract_language_preference(first_name_raw)
            
            # Clean the first name using comprehensive cleaning
            clean_first_name = clean_name_comprehensive(first_name_raw)
            
            # Update the DataFrame
            df.loc[index, 'First Name '] = clean_first_name
            
            # Update language preference if the current one is empty or if we detected from name
            current_lang_pref = row['What is your preferred language for sermons and/or engaging in group discussions?']
            if pd.isna(current_lang_pref) or str(current_lang_pref).strip() == '':
                df.loc[index, 'What is your preferred language for sermons and/or engaging in group discussions?'] = language_pref
                print(f"Updated language preference for {clean_first_name}: {language_pref}")
    
    return df

# Process the main Excel file (contains all participants with correct assignments)
excel_assignments_file = os.path.join(project_root, "input_data", "BADGE_ASSIGNMENTS_FINAL.xlsx")

if os.path.exists(excel_assignments_file):
    print(f"Using Excel assignments file: {excel_assignments_file}")
    df1 = pd.read_excel(excel_assignments_file)
    print(f"Excel assignments file: {len(df1)} rows")
    
    # The Excel file already has the correct format, just need to ensure proper column names
    print("Columns in Excel file:", list(df1.columns))
    
else:
    print(f"Excel assignments file not found: {excel_assignments_file}")
    print("Falling back to CSV file...")
    df1 = pd.read_csv(os.path.join(project_root, "data", "csv", "badge_assignments.csv"))
    print(f"Main CSV file: {len(df1)} rows")
    
    # Convert CSV format to expected format
    def convert_csv_to_expected_format(df):
        """Convert CSV format to the format expected by the badge generation script"""
        # Split full_name into first and last name
        df['First Name '] = df['full_name'].str.split().str[0]
        df['Last Name'] = df['full_name'].str.split().str[1:].str.join(' ')
        
        # Map campus to the expected column name
        df['Which campus ministry are you a part of? In case, you are not part of the listed campus ministries, it\'s open to you as well, and please come and join us! '] = df['campus']
        
        # Map dorm to expected column name
        df['DORM NUMBER'] = df['dorm']
        
        # Map table assignment to expected column name
        df['TABLE #'] = df['table_assignment']
        
        # Map discussion assignment to expected column name
        df['DISCUSSION #'] = df['discussion_assignment']
        
        # Set default language preference
        df['What is your preferred language for sermons and/or engaging in group discussions?'] = 'Amharic'
        
        return df
    
    df1 = convert_csv_to_expected_format(df1)

# Process names with &A and &E indicators in the main file
df1 = process_main_excel_file(df1)
print("Processed main file for language indicators in names")

# All participants are now in the CSV file, so we don't need to process additional files
print(f"Total participants: {len(df1)}")
# Handle both 'role' and 'Role' columns for coordinator count
if 'Role' in df1.columns:
    coordinator_count = len(df1[df1['Role'] == 'COORDINATOR'])
elif 'role' in df1.columns:
    coordinator_count = len(df1[df1['role'] == 'COORDINATOR'])
else:
    coordinator_count = 0
    
print(f"Total coordinators: {coordinator_count}")
print(f"Total regular participants: {len(df1) - coordinator_count}")

# The main data is now df1
df = df1
df_combined = df1

print(f"Final count: {len(df_combined)} participants")

# The CSV file already contains all participants, so we use it as is
# Assign tables and discussions to the combined DataFrame
df = assign_tables_and_discussions(df_combined)

# Note: New Excel file will be created after badge generation with actual assignments

valid_rows = []
skipped_rows = []

for index, row in df.iterrows():
    missing_fields = []
    
    if pd.isna(row['First Name ']) or str(row['First Name ']).strip() == '':
        missing_fields.append('First Name')
    if pd.isna(row['Last Name']) or str(row['Last Name']).strip() == '':
        missing_fields.append('Last Name')
    
    if missing_fields:
        skipped_rows.append(f"Row {index + 2}: Missing {', '.join(missing_fields)}")
    else:
        full_name = f"{str(row['First Name ']).strip().upper()} {str(row['Last Name']).strip().upper()}"
        
        campus_ministry = str(row['Which campus ministry are you a part of? In case, you are not part of the listed campus ministries, it\'s open to you as well, and please come and join us! ']) if not pd.isna(row['Which campus ministry are you a part of? In case, you are not part of the listed campus ministries, it\'s open to you as well, and please come and join us! ']) else 'NEW'
        if 'I don\'t have' in campus_ministry or campus_ministry == 'nan':
            campus_ministry = 'NEW'
        
        # Take only the first ministry if multiple are listed
        if ',' in campus_ministry:
            campus_ministry = campus_ministry.split(',')[0].strip()
        
        campus_location = campus_ministry_locations.get(campus_ministry, campus_ministry)
        
        # Get role from CSV data
        # Handle role column - check both 'role' (from CSV) and 'Role' (from Excel)
        if 'Role' in row.index:
            role = str(row['Role']).strip() if not pd.isna(row['Role']) else 'PARTICIPANT'
        elif 'role' in row.index:
            role = str(row['role']).strip() if not pd.isna(row['role']) else 'PARTICIPANT'
        else:
            role = 'PARTICIPANT'
        
        dorm = str(row['DORM NUMBER']).strip() if not pd.isna(row['DORM NUMBER']) and str(row['DORM NUMBER']).strip() != '' else 'N/A'
        
        # Only assign table/discussion to participants, not coordinators
        if role == 'COORDINATOR':
            table = ''
            discussion = ''
        else:
            table = str(row['TABLE #']).strip() if not pd.isna(row['TABLE #']) else ''
            discussion = str(row['DISCUSSION #']).strip() if not pd.isna(row['DISCUSSION #']) else ''
        
        valid_rows.append({
            'Full_Name': full_name,
            'Campus': campus_location,
            'Role': role,
            'Dorm': dorm,
            'Table': table,
            'Discussion': discussion
        })

print(f"Processing {len(valid_rows)} valid rows")
if skipped_rows:
    print("Skipped rows:")
    for skip in skipped_rows:
        print(f"  {skip}")

template_pres = Presentation(template_file)
template_slide = template_pres.slides[0]

def duplicate_slide(pres, source_slide):
    try:
        blank_slide_layout = pres.slide_layouts[6]
    except:
        blank_slide_layout = pres.slide_layouts[0]
    
    copied_slide = pres.slides.add_slide(blank_slide_layout)
    
    for shape in source_slide.shapes:
        try:
            el = shape.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        except Exception as e:
            print(f"Warning: Could not copy shape: {e}")
            continue
    
    try:
        for rel_id, rel in source_slide.part.rels.items():
            if "notesSlide" not in rel.reltype:
                copied_slide.part.rels.add_relationship(
                    rel.reltype, rel._target, rel_id
                )
    except Exception as e:
        print(f"Warning: Could not copy relationships: {e}")
    
    return copied_slide

output_pres = Presentation()
output_pres.slide_width = template_pres.slide_width
output_pres.slide_height = template_pres.slide_height

# Template only supports 4 people per slide ({{NAME1}} to {{NAME4}})
people_per_slide = 4
num_batches = math.ceil(len(valid_rows) / people_per_slide)

# Track assignments during badge generation
badge_records = []

for batch_idx in range(num_batches):
    start_idx = batch_idx * 4
    end_idx = min(start_idx + 4, len(valid_rows))
    batch_data = valid_rows[start_idx:end_idx]
    
    new_slide = duplicate_slide(output_pres, template_slide)

    # Use existing Excel assignments if available, otherwise generate new ones
    excel_file = os.path.join(project_root, 'input_data', 'BADGE_ASSIGNMENTS_FINAL.xlsx')
    slide_table_assignments = {}
    
    try:
        existing_excel = pd.read_excel(excel_file)
        print(f"Using existing assignments from {excel_file}")
        
        # Load assignments from Excel for this slide
        for i, row_data in enumerate(batch_data):
            badge_num = i + 1
            
            # Find matching person in Excel by name only (ignore slide/slot since order may change)
            person_record = existing_excel[
                (existing_excel['First Name '].astype(str).str.strip() + ' ' + 
                 existing_excel['Last Name'].astype(str).str.strip()) == row_data['Full_Name']
            ]
            
            if not person_record.empty:
                # Use assignment from Excel
                table_assignment = person_record.iloc[0]['TABLE #']
                if pd.isna(table_assignment):
                    table_assignment = ''
                else:
                    table_assignment = str(table_assignment).strip()
                    
                discussion_assignment = person_record.iloc[0]['DISCUSSION #']
                if pd.isna(discussion_assignment):
                    discussion_assignment = ''
                else:
                    discussion_assignment = str(discussion_assignment).strip()
                
                slide_table_assignments[badge_num] = table_assignment
                print(f"  Preserved assignment for {row_data['Full_Name']}: {table_assignment}")
                
                # Record this person's assignments
                badge_record = {
                    'slide': batch_idx + 1,
                    'slot': i,
                    'full_name': row_data['Full_Name'],
                    'role': row_data['Role'],
                    'campus': row_data['Campus'],
                    'dorm': row_data['Dorm'],
                    'table_assignment': table_assignment,
                    'discussion_assignment': discussion_assignment
                }
                badge_records.append(badge_record)
            else:
                # Generate new assignment if not found in CSV
                if row_data['Role'] == 'COORDINATOR':
                    # Coordinators don't get table assignments
                    slide_table_assignments[badge_num] = ''
                    print(f"  Coordinator {row_data['Full_Name']}: No table assignment")
                else:
                    # Generate random assignment for participants
                    import random
                    sat_table_num = random.randint(1, 35)
                    sun_table_num = random.randint(1, 35)
                    sat_table = f"T{sat_table_num}"
                    sun_table = f"T{sun_table_num}"
                    table_text = f"SAT {sat_table} | SUN {sun_table}"
                    
                    slide_table_assignments[badge_num] = table_text
                    print(f"  NEW randomized assignment for {row_data['Full_Name']}: {table_text}")
                
                # Record this person's assignments
                badge_record = {
                    'slide': batch_idx + 1,
                    'slot': i,
                    'full_name': row_data['Full_Name'],
                    'role': row_data['Role'],
                    'campus': row_data['Campus'],
                    'dorm': row_data['Dorm'],
                    'table_assignment': slide_table_assignments[badge_num],
                    'discussion_assignment': row_data['Discussion']
                }
                badge_records.append(badge_record)
                
    except FileNotFoundError:
        print(f"Excel file not found, generating new assignments...")
        
        # Generate new assignments
        for i, row_data in enumerate(batch_data):
            badge_num = i + 1
            
            # Generate unique table assignment for this person
            if row_data['Role'] == 'COORDINATOR':
                # Coordinators don't get table assignments
                table_text = ''
            else:
                # Generate random assignment for participants
                import random
                sat_table_num = random.randint(1, 35)
                sun_table_num = random.randint(1, 35)
                sat_table = f"T{sat_table_num}"
                sun_table = f"T{sun_table_num}"
                table_text = f"SAT {sat_table} | SUN {sun_table}"
            
            # Store this person's unique table assignment
            slide_table_assignments[badge_num] = table_text
            
            # Record this person's assignments
            badge_record = {
                'slide': batch_idx + 1,
                'slot': i,
                'full_name': row_data['Full_Name'],
                'role': row_data['Role'],
                'campus': row_data['Campus'],
                'dorm': row_data['Dorm'],
                'table_assignment': table_text,
                'discussion_assignment': row_data['Discussion']
            }
            badge_records.append(badge_record)
    
    # Process ALL placeholders on the slide in one pass
    for shape in new_slide.shapes:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Process each person's placeholders
                    for i, row_data in enumerate(batch_data):
                        badge_num = i + 1
                        
                        name_placeholder = f"{{{{NAME{badge_num}}}}}"
                        campus_placeholder = f"{{{{CAMPUS{badge_num}}}}}"
                        role_placeholder = f"{{{{ROLE{badge_num}}}}}"
                        dorm_placeholder = f"{{{{DORM{badge_num}}}}}"
                        table_placeholder = f"{{{{TABLE{badge_num}}}}}"
                        discussion_placeholder = f"{{{{DISCUSSION{badge_num}}}}}"
                        
                        if name_placeholder in run.text:
                            run.text = run.text.replace(name_placeholder, row_data['Full_Name'])
                            # Set name color to gold (#FFD700) for better visibility
                            run.font.color.rgb = RGBColor(0xFF, 0xD7, 0x00)
                            # Make name text extra bold
                            run.font.bold = True
                            # Smart sizing based on length - use "MERON BELAYNEH/AMSALU" as max reference (19 chars)
                            name_length = len(row_data['Full_Name'])
                            if name_length > 19:  # Longer than reference name, reduce size
                                if run.font.size:
                                    run.font.size = int(run.font.size * 0.75)
                            elif name_length > 16:  # Slightly long names, small reduction
                                if run.font.size:
                                    run.font.size = int(run.font.size * 0.9)
                            # Names 16 chars or less (like "BETHLEHEM ADUGNA") keep original size
                        elif campus_placeholder in run.text:
                            # For coordinators, show role instead of campus
                            if row_data['Role'] == 'COORDINATOR':
                                display_text = 'COORDINATOR'
                            else:
                                display_text = 'NEW' if row_data['Campus'] == 'NEW' else row_data['Campus']
                            
                            run.text = run.text.replace(campus_placeholder, display_text)
                            if len(display_text) > 25:
                                if run.font.size:
                                    run.font.size = int(run.font.size * 0.85)
                            elif len(display_text) > 20:
                                if run.font.size:
                                    run.font.size = int(run.font.size * 0.9)
                            
                            if "GEORGE WASHINGTON" in display_text or "GEORGE MASON" in display_text:
                                if hasattr(shape, 'left') and shape.left > 100000:
                                    shape.left = max(50000, shape.left - 50000)
                        elif role_placeholder in run.text:
                            # For coordinators, don't show role again (it's already shown in campus field)
                            if row_data['Role'] == 'COORDINATOR':
                                run.text = run.text.replace(role_placeholder, '')
                            else:
                                run.text = run.text.replace(role_placeholder, row_data['Role'])
                        elif dorm_placeholder in run.text:
                            run.text = run.text.replace(dorm_placeholder, row_data['Dorm'])
                        elif table_placeholder in run.text:
                            # Use the unique table assignment for this specific person
                            assignment = slide_table_assignments[badge_num]
                            run.text = run.text.replace(table_placeholder, assignment)
                            # Make table text bold and dark blue to match discussion formatting
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)  # Dark blue (Text 2)
                        elif discussion_placeholder in run.text:
                            if row_data['Discussion'] and row_data['Discussion'] != '':
                                run.text = run.text.replace(discussion_placeholder, row_data['Discussion'])
                            else:
                                run.text = run.text.replace(discussion_placeholder, "")

    # Clear remaining badge placeholders
    remaining_badges = 4 - len(batch_data)
    for shape in new_slide.shapes:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for i in range(remaining_badges):
                        badge_num = len(batch_data) + i + 1
                        
                        name_placeholder = f"{{{{NAME{badge_num}}}}}"
                        campus_placeholder = f"{{{{CAMPUS{badge_num}}}}}"
                        role_placeholder = f"{{{{ROLE{badge_num}}}}}"
                        dorm_placeholder = f"{{{{DORM{badge_num}}}}}"
                        table_placeholder = f"{{{{TABLE{badge_num}}}}}"
                        discussion_placeholder = f"{{{{DISCUSSION{badge_num}}}}}"
                        
                        run.text = run.text.replace(name_placeholder, "")
                        run.text = run.text.replace(campus_placeholder, "")
                        run.text = run.text.replace(role_placeholder, "")
                        run.text = run.text.replace(dorm_placeholder, "")
                        run.text = run.text.replace(table_placeholder, "")
                        run.text = run.text.replace(discussion_placeholder, "")

output_pres.save(output_file)
print(f"Saved {num_batches} slides to {output_file}")

# Use the existing CSV data if it exists, otherwise create new one
csv_file = os.path.join(project_root, 'data', 'csv', 'badge_assignments.csv')
try:
    existing_csv = pd.read_csv(csv_file)
    print(f"\nUsing existing CSV file: {csv_file}")
    
    # Remove duplicates from existing CSV before processing
    print(f"CSV entries before deduplication: {len(existing_csv)}")
    existing_csv = existing_csv.drop_duplicates(subset=['full_name'], keep='first')
    print(f"CSV entries after deduplication: {len(existing_csv)}")
    
    # Create new Excel file using the CSV data
    new_excel_data = []
    for _, record in existing_csv.iterrows():
        # Split the full name back to first and last
        name_parts = record['full_name'].split()
        first_name = name_parts[0] if len(name_parts) > 0 else ''
        last_name = ' '.join(name_parts[1:]) if len(name_parts) > 1 else ''
        
        new_excel_data.append({
            'First Name ': first_name,
            'Last Name': last_name,
            'DORM NUMBER': record['dorm'],
            'TABLE #': record['table_assignment'],
            'DISCUSSION #': record['discussion_assignment'],
            'What is your preferred language for sermons and/or engaging in group discussions?': '',
            'Which campus ministry are you a part of? In case, you are not part of the listed campus ministries, it\'s open to you as well, and please come and join us! ': record['campus'],
            'State you residence in': '',
            'If your are a College Student, which year are you?': '',
            'Gender': ''
        })
    
    # Create DataFrame and save to new Excel file
    new_df = pd.DataFrame(new_excel_data)
    new_excel_filename = os.path.join(project_root, 'input_data', 'BADGE_ASSIGNMENTS_FINAL.xlsx')
    new_df.to_excel(new_excel_filename, index=False)
    print(f"Created new Excel file: {new_excel_filename}")
    
except FileNotFoundError:
    print(f"\nCSV file not found, creating new assignments...")
    
    # Remove duplicates from badge_records before creating final output
    import pandas as pd
    badge_records_df = pd.DataFrame(badge_records)
    print(f"Badge records before deduplication: {len(badge_records_df)}")
    
    # Remove duplicates by full name, keeping the first occurrence
    badge_records_df = badge_records_df.drop_duplicates(subset=['full_name'], keep='first')
    print(f"Badge records after deduplication: {len(badge_records_df)}")
    
    badge_records = badge_records_df.to_dict('records')
    
    # Create new Excel file with the exact data from badges
    new_excel_data = []
    for record in badge_records:
        # Split the full name back to first and last
        name_parts = record['full_name'].split()
        first_name = name_parts[0] if len(name_parts) > 0 else ''
        last_name = ' '.join(name_parts[1:]) if len(name_parts) > 1 else ''
        
        new_excel_data.append({
            'First Name ': first_name,
            'Last Name': last_name,
            'DORM NUMBER': record['dorm'],
            'TABLE #': record['table_assignment'],
            'DISCUSSION #': record['discussion_assignment'],
            'What is your preferred language for sermons and/or engaging in group discussions?': '',
            'Which campus ministry are you a part of? In case, you are not part of the listed campus ministries, it\'s open to you as well, and please come and join us! ': record['campus'],
            'State you residence in': '',
            'If your are a College Student, which year are you?': '',
            'Gender': ''
        })
    
    # Create DataFrame and save to new Excel file
    new_df = pd.DataFrame(new_excel_data)
    new_excel_filename = os.path.join(project_root, 'input_data', 'BADGE_ASSIGNMENTS_FINAL.xlsx')
    new_df.to_excel(new_excel_filename, index=False)
    print(f"Created new Excel file: {new_excel_filename}")
    
    # Also save badge assignments to CSV for reference
    badge_assignments_df = pd.DataFrame(badge_records)
    badge_assignments_df.to_csv(os.path.join(project_root, 'input_data', 'badge_assignments.csv'), index=False)
    print(f"Saved badge assignments to {os.path.join(project_root, 'input_data', 'badge_assignments.csv')}")

print(f"\nBadge generation complete!")
print(f"Files created:")
print(f"  - {output_file} (PowerPoint with badges)")
print(f"  - {new_excel_filename} (Excel matching badge order)")
print(f"  - {os.path.join(project_root, 'input_data', 'badge_assignments.csv')} (CSV for reference)")