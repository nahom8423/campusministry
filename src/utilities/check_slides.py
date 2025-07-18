#!/usr/bin/env python3

from pptx import Presentation
import os

def check_slides():
    pptx_path = "/Users/nahomnigatu/Downloads/campusministrybadges/filled_badges.pptx"
    
    if not os.path.exists(pptx_path):
        print(f"Error: {pptx_path} not found")
        return
    
    try:
        prs = Presentation(pptx_path)
        
        # Check first 3 slides
        for slide_num in range(1, min(4, len(prs.slides) + 1)):
            slide_index = slide_num - 1
            slide = prs.slides[slide_index]
            
            print(f"\n=== SLIDE {slide_num} ===")
            
            # Look for all text in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if text:
                        print(f"Text found: {text}")
                
                # Check for tables
                if hasattr(shape, "table"):
                    print("Table found:")
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            if cell.text.strip():
                                print(f"  Row {row_idx}, Col {col_idx}: {cell.text.strip()}")
                        
    except Exception as e:
        print(f"Error reading presentation: {e}")

if __name__ == "__main__":
    check_slides()