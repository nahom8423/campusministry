#!/usr/bin/env python3

from pptx import Presentation
import os

def detailed_check():
    pptx_path = "/Users/nahomnigatu/Downloads/campusministrybadges/filled_badges.pptx"
    
    if not os.path.exists(pptx_path):
        print(f"Error: {pptx_path} not found")
        return
    
    try:
        prs = Presentation(pptx_path)
        
        # Check first 3 slides in detail
        for slide_num in range(1, min(4, len(prs.slides) + 1)):
            slide_index = slide_num - 1
            slide = prs.slides[slide_index]
            
            print(f"\n=== SLIDE {slide_num} DETAILED ANALYSIS ===")
            
            # Look for each shape with coordinates
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    left = shape.left
                    top = shape.top
                    print(f"Shape {shape_idx}: '{text}' at position ({left}, {top})")
                    
                    # If it's a table assignment, try to identify which badge it belongs to
                    if "SAT T" in text or "SUN T" in text:
                        print(f"  ** TABLE ASSIGNMENT FOUND: {text} **")
                
                # Check for tables
                if hasattr(shape, "table"):
                    print(f"Shape {shape_idx}: TABLE")
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            if cell.text.strip():
                                print(f"  Table Cell [{row_idx}][{col_idx}]: '{cell.text.strip()}'")
                                
    except Exception as e:
        print(f"Error reading presentation: {e}")

if __name__ == "__main__":
    detailed_check()