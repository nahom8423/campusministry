import pandas as pd
from pptx import Presentation
import copy

excel_file = "General Assembly Dorm Assignment.xlsx"
template_file = "CAMPUSMINISTRYBADGE.pptx"
output_file = "filled_badges.pptx"

df = pd.read_excel(excel_file)

valid_rows = []
skipped_rows = []

for index, row in df.iterrows():
    if index >= 137:
        break
        
    missing_fields = []
    
    if pd.isna(row['First Name']) or str(row['First Name']).strip() == '':
        missing_fields.append('First Name')
    if pd.isna(row['Last Name']) or str(row['Last Name']).strip() == '':
        missing_fields.append('Last Name')
    
    attendee_type = str(row['Attendee Type']).strip().upper() if not pd.isna(row['Attendee Type']) else ''
    
    if attendee_type != 'PARTICIPANT':
        if attendee_type == '':
            missing_fields.append('Attendee Type')
        if pd.isna(row['Room']) or str(row['Room']).strip() == '':
            missing_fields.append('Room')
    else:
        if pd.isna(row['Room']) or str(row['Room']).strip() == '':
            missing_fields.append('Room')
    
    if missing_fields:
        skipped_rows.append(f"Row {index + 2}: Missing {', '.join(missing_fields)}")
    else:
        full_name = f"{str(row['First Name']).strip().upper()} {str(row['Last Name']).strip().upper()}"
        
        campus = str(row['Campus Minstry']) if not pd.isna(row['Campus Minstry']) else 'NEW'
        if 'I don\'t have' in campus or campus == 'nan':
            campus = 'NEW'
        
        role = attendee_type
        dorm = str(row['Room']).strip()
        
        valid_rows.append({
            'Full_Name': full_name,
            'Campus': campus,
            'Role': role,
            'Dorm': dorm
        })

print(f"Processing up to row 137: {len(valid_rows)} valid rows")
if skipped_rows:
    print("Skipped rows:")
    for skip in skipped_rows:
        print(f"  {skip}")

template_pres = Presentation(template_file)
template_slide = template_pres.slides[0]

def duplicate_slide(pres, source_slide):
    try:
        blank_slide_layout = pres.slide_layouts[6]
    except:
        blank_slide_layout = pres.slide_layouts[0]
    
    copied_slide = pres.slides.add_slide(blank_slide_layout)
    
    for shape in source_slide.shapes:
        try:
            el = shape.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        except Exception as e:
            print(f"Warning: Could not copy shape: {e}")
            continue
    
    try:
        for rel_id, rel in source_slide.part.rels.items():
            if "notesSlide" not in rel.reltype:
                copied_slide.part.rels.add_relationship(
                    rel.reltype, rel._target, rel_id
                )
    except Exception as e:
        print(f"Warning: Could not copy relationships: {e}")
    
    return copied_slide

output_pres = Presentation()
output_pres.slide_width = template_pres.slide_width
output_pres.slide_height = template_pres.slide_height

import math

num_batches = math.ceil(len(valid_rows) / 4)

for batch_idx in range(num_batches):
    start_idx = batch_idx * 4
    end_idx = min(start_idx + 4, len(valid_rows))
    batch_data = valid_rows[start_idx:end_idx]
    
    new_slide = duplicate_slide(output_pres, template_slide)

    for i, row_data in enumerate(batch_data):
        badge_num = i + 1
        
        name_placeholder = f"{{{{NAME{badge_num}}}}}"
        campus_placeholder = f"{{{{CAMPUS{badge_num}}}}}"
        role_placeholder = f"{{{{ROLE{badge_num}}}}}"
        dorm_placeholder = f"{{{{DORM{badge_num}}}}}"
        
        for shape in new_slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if name_placeholder in run.text:
                            run.text = run.text.replace(name_placeholder, row_data['Full_Name'])
                        if campus_placeholder in run.text:
                            run.text = run.text.replace(campus_placeholder, row_data['Campus'])
                        if role_placeholder in run.text:
                            run.text = run.text.replace(role_placeholder, row_data['Role'])
                        if dorm_placeholder in run.text:
                            run.text = run.text.replace(dorm_placeholder, row_data['Dorm'])

    remaining_badges = 4 - len(batch_data)
    for i in range(remaining_badges):
        badge_num = len(batch_data) + i + 1
        
        name_placeholder = f"{{{{NAME{badge_num}}}}}"
        campus_placeholder = f"{{{{CAMPUS{badge_num}}}}}"
        role_placeholder = f"{{{{ROLE{badge_num}}}}}"
        dorm_placeholder = f"{{{{DORM{badge_num}}}}}"
        
        for shape in new_slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace(name_placeholder, "")
                        run.text = run.text.replace(campus_placeholder, "")
                        run.text = run.text.replace(role_placeholder, "")
                        run.text = run.text.replace(dorm_placeholder, "")

output_pres.save(output_file)
print(f"Saved {num_batches} slides to {output_file}")