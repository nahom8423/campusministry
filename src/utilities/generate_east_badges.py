import pandas as pd
from pptx import Presentation
import copy
import math

excel_file = "East-registration .xlsx"
template_file = "mkdcbadges.pptx"
output_file = "filled_badges.pptx"

city_corrections = {
    "DC": "WASHINGTON D.C.",
    "Dc": "WASHINGTON D.C.", 
    "Virginia": "VIRGINIA",
    "Verginia": "VIRGINIA",
    "Verginia ": "VIRGINIA",
    "Virginia ": "VIRGINIA",
    "Verginia North": "NORTHERN VIRGINIA",
    "Boston": "BOSTON",
    "Bosten": "BOSTON",
    "Bosten ": "BOSTON",
    "New York": "NEW YORK",
    "Newyork": "NEW YORK",
    "Newyork ": "NEW YORK",
    "Carolina": "NORTH CAROLINA",
    " N Carolina ": "NORTH CAROLINA",
    "Indiana": "INDIANA",
    "Indiana ": "INDIANA",
    "Ohio": "OHIO",
    "Ohio ": "OHIO",
    "Ve": "VIRGINIA",
    "Ver": "VIRGINIA",
    "Guest": "GUEST"
}

df = pd.read_excel(excel_file)

valid_rows = []
skipped_rows = []

for index, row in df.iterrows():
    missing_fields = []
    
    name_col = df.columns[0]
    city_col = df.columns[1]
    
    if pd.isna(row[name_col]) or str(row[name_col]).strip() == '':
        missing_fields.append('Name')
    
    if missing_fields:
        skipped_rows.append(f"Row {index + 2}: Missing {', '.join(missing_fields)}")
    else:
        full_name = str(row[name_col]).strip().upper()
        
        city_raw = str(row[city_col]) if not pd.isna(row[city_col]) else ''
        city_corrected = city_corrections.get(city_raw, city_raw.upper())
        
        badge_text = f"{full_name},\n{city_corrected}"
        
        valid_rows.append({
            'Badge_Text': badge_text
        })

print(f"Processing {len(valid_rows)} valid rows")
if skipped_rows:
    print("Skipped rows:")
    for skip in skipped_rows:
        print(f"  {skip}")

template_pres = Presentation(template_file)
template_slide = template_pres.slides[0]

def duplicate_slide(pres, source_slide):
    try:
        blank_slide_layout = pres.slide_layouts[6]
    except:
        blank_slide_layout = pres.slide_layouts[0]
    
    copied_slide = pres.slides.add_slide(blank_slide_layout)
    
    for shape in source_slide.shapes:
        try:
            el = shape.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        except Exception as e:
            print(f"Warning: Could not copy shape: {e}")
            continue
    
    try:
        for rel_id, rel in source_slide.part.rels.items():
            if "notesSlide" not in rel.reltype:
                copied_slide.part.rels.add_relationship(
                    rel.reltype, rel._target, rel_id
                )
    except Exception as e:
        print(f"Warning: Could not copy relationships: {e}")
    
    return copied_slide

output_pres = Presentation()
output_pres.slide_width = template_pres.slide_width
output_pres.slide_height = template_pres.slide_height

num_batches = math.ceil(len(valid_rows) / 4)

for batch_idx in range(num_batches):
    start_idx = batch_idx * 4
    end_idx = min(start_idx + 4, len(valid_rows))
    batch_data = valid_rows[start_idx:end_idx]
    
    new_slide = duplicate_slide(output_pres, template_slide)

    for i, row_data in enumerate(batch_data):
        badge_num = i + 1
        text_placeholder = f"{{{{Text{badge_num}}}}}"
        
        for shape in new_slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if text_placeholder in run.text:
                            run.text = run.text.replace(text_placeholder, row_data['Badge_Text'])
                            if len(row_data['Badge_Text']) > 25:
                                if run.font.size:
                                    run.font.size = int(run.font.size * 0.7)
                            elif len(row_data['Badge_Text']) > 20:
                                if run.font.size:
                                    run.font.size = int(run.font.size * 0.8)

    remaining_badges = 4 - len(batch_data)
    for i in range(remaining_badges):
        badge_num = len(batch_data) + i + 1
        text_placeholder = f"{{{{Text{badge_num}}}}}"
        
        for shape in new_slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace(text_placeholder, "")

output_pres.save(output_file)
print(f"Saved {num_batches} slides to {output_file}")