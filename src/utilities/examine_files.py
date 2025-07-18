import pandas as pd
import os
from pptx import Presentation

# Read the Excel file
excel_path = '/Users/nahomnigatu/Downloads/campusministrybadges/BADGE_ASSIGNMENTS_FINAL.xlsx'
df = pd.read_excel(excel_path)

print("=" * 60)
print("EXCEL FILE ANALYSIS")
print("=" * 60)
print(f"Shape: {df.shape}")
print(f"Columns: {df.columns.tolist()}")
print("\nFirst 10 rows:")
print(df.head(10))

print("\nTable assignment distribution:")
if 'Table' in df.columns:
    print(df['Table'].value_counts().sort_index())

print("\n" + "=" * 60)
print("POWERPOINT FILE ANALYSIS")
print("=" * 60)

# Read the PowerPoint file
pptx_path = '/Users/nahomnigatu/Downloads/campusministrybadges/filled_badges.pptx'
prs = Presentation(pptx_path)

print(f"Total slides: {len(prs.slides)}")

# Examine first slide
if len(prs.slides) > 0:
    slide = prs.slides[0]
    print(f"\nFirst slide analysis:")
    print(f"Number of shapes: {len(slide.shapes)}")
    
    # Look for text in shapes
    for i, shape in enumerate(slide.shapes):
        if hasattr(shape, 'text') and shape.text.strip():
            print(f"Shape {i}: {shape.text.strip()}")
        elif hasattr(shape, 'text_frame') and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        print(f"Shape {i} text: {run.text.strip()}")

# Let's also check a few more slides to see the pattern
print(f"\nChecking first 3 slides for table assignments:")
for slide_num in range(min(3, len(prs.slides))):
    slide = prs.slides[slide_num]
    print(f"\nSlide {slide_num + 1}:")
    
    # Collect all text from the slide
    all_text = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            all_text.append(shape.text.strip())
        elif hasattr(shape, 'text_frame') and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        all_text.append(run.text.strip())
    
    # Filter for table assignments (assuming they contain "Table")
    table_assignments = [text for text in all_text if 'Table' in text]
    print(f"Table assignments found: {table_assignments}")
    
    # Also look for names (assuming they don't contain "Table" but are substantial text)
    names = [text for text in all_text if 'Table' not in text and len(text) > 3 and not text.isdigit()]
    print(f"Names found: {names}")