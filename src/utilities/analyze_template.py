#!/usr/bin/env python3
"""
Script to analyze PowerPoint template and extract all placeholders
"""

import re
from pptx import Presentation

def analyze_template(template_path):
    """
    Analyze PowerPoint template to find all placeholders
    """
    try:
        pres = Presentation(template_path)
        print(f"Analyzing template: {template_path}")
        print(f"Number of slides: {len(pres.slides)}")
        
        # Analyze first slide (template slide)
        template_slide = pres.slides[0]
        print(f"Number of shapes on template slide: {len(template_slide.shapes)}")
        
        # Find all placeholders
        placeholders = set()
        
        for shape_idx, shape in enumerate(template_slide.shapes):
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Find placeholders in format {{PLACEHOLDER}}
                        matches = re.findall(r'\{\{([^}]+)\}\}', run.text)
                        for match in matches:
                            placeholders.add(match)
                        
                        # Print text content for debugging
                        if run.text.strip():
                            print(f"Shape {shape_idx} text: {repr(run.text)}")
        
        print(f"\nFound placeholders:")
        sorted_placeholders = sorted(list(placeholders))
        for placeholder in sorted_placeholders:
            print(f"  {{{{{placeholder}}}}}")
        
        # Check for specific patterns
        name_placeholders = [p for p in sorted_placeholders if p.startswith('NAME')]
        campus_placeholders = [p for p in sorted_placeholders if p.startswith('CAMPUS')]
        role_placeholders = [p for p in sorted_placeholders if p.startswith('ROLE')]
        dorm_placeholders = [p for p in sorted_placeholders if p.startswith('DORM')]
        table_placeholders = [p for p in sorted_placeholders if p.startswith('TABLE')]
        discussion_placeholders = [p for p in sorted_placeholders if p.startswith('DISCUSSION')]
        
        print(f"\nPlaceholder analysis:")
        print(f"NAME placeholders: {name_placeholders}")
        print(f"CAMPUS placeholders: {campus_placeholders}")
        print(f"ROLE placeholders: {role_placeholders}")
        print(f"DORM placeholders: {dorm_placeholders}")
        print(f"TABLE placeholders: {table_placeholders}")
        print(f"DISCUSSION placeholders: {discussion_placeholders}")
        
        # Determine maximum number of people per slide
        max_people = 0
        if name_placeholders:
            # Extract numbers from NAME placeholders
            name_numbers = []
            for name_placeholder in name_placeholders:
                match = re.search(r'NAME(\d+)', name_placeholder)
                if match:
                    name_numbers.append(int(match.group(1)))
            if name_numbers:
                max_people = max(name_numbers)
        
        print(f"\nMaximum people per slide supported: {max_people}")
        
        # Check if 6 people per slide is supported
        required_for_6_people = [
            'NAME1', 'NAME2', 'NAME3', 'NAME4', 'NAME5', 'NAME6',
            'CAMPUS1', 'CAMPUS2', 'CAMPUS3', 'CAMPUS4', 'CAMPUS5', 'CAMPUS6',
            'ROLE1', 'ROLE2', 'ROLE3', 'ROLE4', 'ROLE5', 'ROLE6',
            'DORM1', 'DORM2', 'DORM3', 'DORM4', 'DORM5', 'DORM6',
            'TABLE1', 'TABLE2', 'TABLE3', 'TABLE4', 'TABLE5', 'TABLE6',
            'DISCUSSION1', 'DISCUSSION2', 'DISCUSSION3', 'DISCUSSION4', 'DISCUSSION5', 'DISCUSSION6'
        ]
        
        missing_for_6_people = []
        for required in required_for_6_people:
            if required not in sorted_placeholders:
                missing_for_6_people.append(required)
        
        if missing_for_6_people:
            print(f"\nMissing placeholders for 6 people per slide:")
            for missing in missing_for_6_people:
                print(f"  {{{{{missing}}}}}")
            print(f"\nRecommendation: Keep 4 people per slide (template supports up to {max_people} people)")
        else:
            print(f"\nTemplate supports 6 people per slide!")
        
        return {
            'max_people': max_people,
            'placeholders': sorted_placeholders,
            'supports_6_people': len(missing_for_6_people) == 0,
            'missing_for_6_people': missing_for_6_people
        }
        
    except Exception as e:
        print(f"Error analyzing template: {e}")
        return None

if __name__ == "__main__":
    template_file = "CAMPUSMINISTRYBADGE.pptx"
    result = analyze_template(template_file)
    
    if result:
        print(f"\n{'='*60}")
        print("SUMMARY:")
        print(f"Template supports up to {result['max_people']} people per slide")
        if result['supports_6_people']:
            print("✓ Template supports 6 people per slide")
        else:
            print("✗ Template does not support 6 people per slide")
            print(f"  Missing {len(result['missing_for_6_people'])} placeholders")
            print("  Recommendation: Keep 4 people per slide")
        print(f"{'='*60}")